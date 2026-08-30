"""Generate BeaconCall's five-slide, Guava-focused hackathon deck."""

from __future__ import annotations

import subprocess
from html import escape
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "artifacts"
PPTX_OUT = OUT_DIR / "BeaconCall-Guava-Demo.pptx"
HTML_OUT = OUT_DIR / "BeaconCall-Guava-Demo.html"
PDF_OUT = OUT_DIR / "BeaconCall-Guava-Demo.pdf"

W, H = 13.333, 7.5
BG, PANEL, TEXT = "080B09", "0D110E", "EDF4EC"
MUTED, LINE, GREEN = "899188", "29302A", "62D84E"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, *, fill=BG, line_color=None, rounded=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line_color or fill)
    return shape


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=TEXT,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = value
    run.font.name = {"Aptos": "Helvetica Neue", "Aptos Mono": "Menlo"}.get(font, font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def line(slide, x1, y1, x2, y2, *, color=LINE, width=1):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    return connector


def base_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)
    return slide


def eyebrow(slide, label: str, number: str):
    text(slide, label, 0.65, 0.4, 6.5, 0.3, size=10, color=GREEN, bold=True, font="Aptos Mono")
    text(slide, number, 12.2, 0.4, 0.5, 0.3, size=10, color=MUTED, font="Aptos Mono", align=PP_ALIGN.RIGHT)


def footer(slide, label="BEACONCALL / GUAVA INBOUND VOICE"):
    line(slide, 0.65, 7.1, 12.68, 7.1)
    text(slide, label, 0.65, 7.18, 8, 0.16, size=7, color=MUTED, font="Aptos Mono")


def person(slide, x, y, scale=1.0):
    head = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.47 * scale),
        Inches(y),
        Inches(0.34 * scale),
        Inches(0.34 * scale),
    )
    head.fill.solid()
    head.fill.fore_color.rgb = rgb(TEXT)
    head.line.color.rgb = rgb(TEXT)
    rect(slide, x + 0.4 * scale, y + 0.38 * scale, 0.48 * scale, 0.92 * scale, fill=TEXT, rounded=True)
    rect(slide, x + 0.22 * scale, y + 0.44 * scale, 0.18 * scale, 0.78 * scale, fill=TEXT, rounded=True)
    rect(slide, x + 0.88 * scale, y + 0.44 * scale, 0.18 * scale, 0.78 * scale, fill=TEXT, rounded=True)
    rect(slide, x + 0.42 * scale, y + 1.25 * scale, 0.18 * scale, 0.9 * scale, fill=TEXT, rounded=True)
    rect(slide, x + 0.68 * scale, y + 1.25 * scale, 0.18 * scale, 0.9 * scale, fill=TEXT, rounded=True)


def add_title(prs: Presentation):
    slide = base_slide(prs)
    rect(slide, 0, 0, 0.16, H, fill=GREEN)
    text(slide, "HACKATHON DEMO", 0.72, 0.58, 4, 0.3, size=10, color=GREEN, bold=True, font="Aptos Mono")
    text(slide, "BEACONCALL", 0.68, 1.28, 7.1, 1.0, size=56, bold=True)
    text(slide, "CAMERA SIGHTING.\nGUAVA BRIEFING.\nHUMAN RESPONSE.", 0.73, 2.42, 5.2, 1.7, size=22, color=MUTED, bold=True)
    rect(slide, 0.73, 5.65, 4.15, 0.46, fill=GREEN)
    text(slide, "MAC NOW  /  UNITREE G1 LATER", 0.86, 5.77, 3.9, 0.2, size=9, color=BG, bold=True, font="Aptos Mono")
    rect(slide, 7.15, 0.72, 5.55, 5.95, fill=PANEL, line_color=LINE)
    text(slide, "ROBOT CAMERA / PRESENCE ONLY", 7.48, 1.02, 4.8, 0.28, size=9, color=MUTED, font="Aptos Mono")
    rect(slide, 8.25, 1.58, 3.35, 4.2, fill="111713", line_color=GREEN)
    line(slide, 8.25, 3.68, 11.6, 3.68, color=GREEN, width=1.2)
    person(slide, 9.3, 2.05)
    text(slide, "PERSON  91%", 8.4, 5.42, 2.3, 0.25, size=10, color=GREEN, bold=True, font="Aptos Mono")
    text(slide, "guava.", 11.25, 6.82, 1.2, 0.32, size=17, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)


def add_why(prs: Presentation):
    slide = base_slide(prs)
    eyebrow(slide, "WHY WE BUILT IT", "02")
    text(slide, "A visual alert is\neasy to miss.", 0.65, 1.05, 5.8, 1.35, size=42, bold=True)
    text(slide, "A conversation\ngets a response.", 6.72, 1.05, 5.8, 1.35, size=42, color=GREEN, bold=True)
    line(slide, 0.65, 2.78, 12.68, 2.78)
    reasons = [
        ("SEE", "Confirm a person across three frames."),
        ("SPEAK", "Turn camera metadata into a clear briefing."),
        ("RESPOND", "Record what the operator decides."),
    ]
    for index, (label, detail) in enumerate(reasons):
        x = 0.65 + index * 4.08
        text(slide, f"0{index + 1}", x, 3.4, 0.45, 0.3, size=10, color=GREEN, font="Aptos Mono")
        text(slide, label, x, 3.88, 3.4, 0.4, size=22, bold=True)
        text(slide, detail, x, 4.55, 3.25, 0.85, size=14, color=MUTED)
    text(slide, "NO FALL ACTING. NO MEDICAL CLAIM. PRESENCE ONLY.", 0.65, 6.38, 8, 0.3, size=9, color=GREEN, bold=True, font="Aptos Mono")
    footer(slide)


def add_demo(prs: Presentation):
    slide = base_slide(prs)
    eyebrow(slide, "LIVE DEMO / 45–60 SECONDS", "03")
    text(slide, "ONE PERSON. ONE CALL.\nONE CLEAR HANDOFF.", 0.65, 0.98, 8.5, 1.15, size=38, bold=True)
    steps = [
        ("01", "EMPTY FRAME", "Show live camera scanning."),
        ("02", "WALK IN", "Hold for three detections."),
        ("03", "CALL GUAVA", "Ask what the camera saw."),
        ("04", "ACKNOWLEDGE", "Choose inspect or monitor."),
    ]
    for index, (number, label, detail) in enumerate(steps):
        x = 0.65 + index * 3.03
        rect(slide, x, 2.64, 2.62, 2.35, fill=PANEL, line_color=GREEN if index == 2 else LINE)
        text(slide, number, x + 0.22, 2.84, 0.4, 0.25, size=9, color=GREEN, bold=True, font="Aptos Mono")
        text(slide, label, x + 0.22, 3.39, 2.15, 0.38, size=18, bold=True)
        text(slide, detail, x + 0.22, 4.06, 2.05, 0.62, size=12, color=MUTED)
        if index < 3:
            text(slide, "→", x + 2.66, 3.59, 0.3, 0.4, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    text(slide, 'ASK: “WHAT DID THE ROBOT CAMERA SEE?”', 0.65, 5.75, 7, 0.35, size=15, color=GREEN, bold=True, font="Aptos Mono")
    footer(slide)


def add_architecture(prs: Presentation):
    slide = base_slide(prs)
    eyebrow(slide, "SYSTEM DESIGN", "04")
    text(slide, "CONTEXT MOVES.\nTHE BOUNDARY STAYS.", 0.65, 0.92, 6.5, 1.1, size=37, bold=True)
    nodes = [
        ("MAC CAMERA", "G1 stream later"),
        ("COCO-SSD", "person presence"),
        ("INCIDENT", "local JSON"),
        ("GUAVA", "inbound Expert"),
        ("RESPONDER", "voice decision"),
    ]
    for index, (label, detail) in enumerate(nodes):
        x = 0.65 + index * 2.48
        rect(slide, x, 2.62, 2.05, 1.25, fill=PANEL, line_color=GREEN if label == "GUAVA" else LINE)
        text(slide, label, x + 0.17, 2.91, 1.72, 0.26, size=12, color=GREEN if label == "GUAVA" else TEXT, bold=True, font="Aptos Mono")
        text(slide, detail, x + 0.17, 3.35, 1.72, 0.2, size=9, color=MUTED)
        if index < 4:
            text(slide, "→", x + 2.06, 3.03, 0.37, 0.3, size=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    rect(slide, 0.65, 4.58, 12.02, 1.25, fill="0A0E0B", line_color=LINE)
    text(slide, "PRIVACY BOUNDARY", 0.9, 4.89, 2.25, 0.3, size=10, color=GREEN, bold=True, font="Aptos Mono")
    text(slide, "IMAGE STAYS LOCAL", 3.4, 4.84, 2.6, 0.35, size=17, bold=True)
    text(slide, "Guava receives only time, region, confidence, and the presence-only limitation.", 6.15, 4.82, 5.9, 0.55, size=12, color=MUTED)
    text(slide, "SWAP MAC CAMERA → G1 CAMERA. KEEP THE GUAVA EXPERT.", 0.65, 6.38, 8.5, 0.3, size=9, color=GREEN, bold=True, font="Aptos Mono")
    footer(slide)


def add_guava(prs: Presentation):
    slide = base_slide(prs)
    eyebrow(slide, "SPONSOR INTEGRATION", "05")
    text(slide, "guava.", 0.65, 0.98, 4.8, 0.85, size=54, color=GREEN, bold=True)
    text(slide, "VOICE IS THE HANDOFF.", 0.65, 2.0, 6.3, 0.65, size=31, bold=True)
    text(slide, "Not a logo. The live interaction layer.", 0.65, 2.75, 5.1, 0.4, size=15, color=MUTED)
    facts = [
        ("DIALOG SYSTEM", "Audio · speech recognition · conversation · voice"),
        ("LOCAL EXPERT", "Loads the newest verified camera sighting"),
        ("STRUCTURED TASK", "Briefs, asks, constrains, and records response"),
        ("INBOUND PHONE", "Works from a Mac without a public webhook"),
    ]
    for index, (label, detail) in enumerate(facts):
        y = 0.95 + index * 1.24
        text(slide, f"0{index + 1}", 7.02, y, 0.35, 0.24, size=9, color=GREEN, font="Aptos Mono")
        text(slide, label, 7.62, y - 0.02, 4.2, 0.3, size=14, bold=True, font="Aptos Mono")
        text(slide, detail, 7.62, y + 0.42, 4.55, 0.42, size=11, color=MUTED)
        if index < 3:
            line(slide, 7.62, y + 0.94, 12.48, y + 0.94)
    rect(slide, 0.65, 5.45, 5.55, 0.67, fill=GREEN)
    text(slide, "PROOF: LIVE CALL + GUAVA CONVERSATIONS", 0.86, 5.67, 5.14, 0.22, size=9, color=BG, bold=True, font="Aptos Mono")
    text(slide, "goguava.ai/docs/inbound-form-filling", 7.62, 6.45, 4.9, 0.22, size=8, color=MUTED, font="Aptos Mono", align=PP_ALIGN.RIGHT)
    footer(slide, "BEACONCALL / BUILT AROUND GUAVA INBOUND")


def write_html() -> None:
    slides = [
        ("HACKATHON DEMO", "BEACONCALL", "CAMERA SIGHTING · GUAVA BRIEFING · HUMAN RESPONSE"),
        ("WHY", "A visual alert is easy to miss. A conversation gets a response.", "SEE · SPEAK · RESPOND"),
        ("LIVE DEMO", "One person. One call. One clear handoff.", "EMPTY FRAME → WALK IN → CALL GUAVA → ACKNOWLEDGE"),
        ("ARCHITECTURE", "Context moves. The boundary stays.", "CAMERA → COCO-SSD → INCIDENT → GUAVA → RESPONDER"),
        ("guava.", "Voice is the handoff.", "DIALOG SYSTEM · LOCAL EXPERT · STRUCTURED TASK · INBOUND PHONE"),
    ]
    sections = "\n".join(
        f'<section><div class="n">0{i}</div><div class="k">{escape(kicker)}</div><h1>{escape(title)}</h1><p>{escape(body)}</p><div class="g">guava.</div></section>'
        for i, (kicker, title, body) in enumerate(slides, 1)
    )
    HTML_OUT.write_text(
        f"""<!doctype html><html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width"><title>BeaconCall — Guava Demo</title><style>
*{{box-sizing:border-box}}html{{scroll-snap-type:y mandatory}}body{{margin:0;background:#080b09;color:#edf4ec;font-family:Arial,sans-serif}}section{{height:100vh;scroll-snap-align:start;padding:8vh 7vw;position:relative;border-left:12px solid #62d84e;display:flex;flex-direction:column;justify-content:center}}.n{{position:absolute;right:5vw;top:5vh;color:#899188;font:14px monospace}}.k{{color:#62d84e;font:700 14px monospace;letter-spacing:.18em;margin-bottom:4vh}}h1{{font-size:clamp(48px,5.7vw,92px);letter-spacing:-.055em;line-height:.9;max-width:18ch;margin:0}}p{{font:700 clamp(14px,1.7vw,24px) monospace;color:#899188;letter-spacing:.08em;margin-top:8vh;max-width:42ch}}.g{{position:absolute;right:5vw;bottom:5vh;color:#62d84e;font-size:30px;font-weight:900}}@media(max-width:700px){{section{{padding:7vh 8vw}}h1{{font-size:12vw}}p{{font-size:2.7vw}}}}@page{{size:13.333in 7.5in;margin:0}}@media print{{html,body{{width:13.333in;background:#080b09}}section{{width:13.333in;height:7.5in;break-after:page;page-break-after:always}}section:last-child{{break-after:auto;page-break-after:auto}}}}
</style></head><body>{sections}</body></html>"""
    )


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    add_title(prs)
    add_why(prs)
    add_demo(prs)
    add_architecture(prs)
    add_guava(prs)
    prs.save(PPTX_OUT)
    write_html()
    chrome = Path("/Applications/Google Chrome.app/Contents/MacOS/Google Chrome")
    if chrome.exists():
        subprocess.run(
            [
                str(chrome),
                "--headless=new",
                "--disable-gpu",
                "--no-pdf-header-footer",
                f"--print-to-pdf={PDF_OUT}",
                HTML_OUT.as_uri(),
            ],
            check=True,
        )
    print(PPTX_OUT)
    print(HTML_OUT)
    if PDF_OUT.exists():
        print(PDF_OUT)


if __name__ == "__main__":
    main()
